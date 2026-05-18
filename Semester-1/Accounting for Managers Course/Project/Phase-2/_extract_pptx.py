from pptx import Presentation

pptx_path = r'C:\Personal\Personal Docs\2025-Doc\Daily\MBA\Semester-1\Accounting for Managers Course\Project\Phase-2\Financial statement analysis.pptx'
prs = Presentation(pptx_path)

for i, slide in enumerate(prs.slides, 1):
    print(f'--- Slide {i} ---')
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            print(shape.text.strip())
        if shape.has_table:
            table = shape.table
            for row_idx, row in enumerate(table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                print(' | '.join(cells))
    print()
